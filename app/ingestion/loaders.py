from pathlib import Path
from typing import Iterable
import pandas as pd
from docx import Document
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".txt", ".md"}


def load_docx(path: Path) -> list[str]:
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return parts


def load_pptx(path: Path) -> list[str]:
    prs = Presentation(path)
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"Slide {idx}: " + "\n".join(slide_text))
    return parts


def load_xlsx(path: Path) -> list[str]:
    parts = []
    excel = pd.ExcelFile(path)
    for sheet in excel.sheet_names:
        df = excel.parse(sheet).fillna("")
        if df.empty:
            continue
        headers = [str(c) for c in df.columns]
        for i, row in df.iterrows():
            values = []
            for header, value in zip(headers, row.tolist()):
                value = str(value).strip()
                if value:
                    values.append(f"{header}: {value}")
            if values:
                parts.append(f"Sheet: {sheet} | Row {i + 1} | " + " | ".join(values))
    return parts


def load_text(path: Path) -> list[str]:
    return [path.read_text(encoding="utf-8", errors="ignore")]


def load_file(path: Path) -> list[str]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return load_docx(path)
    if suffix == ".pptx":
        return load_pptx(path)
    if suffix == ".xlsx":
        return load_xlsx(path)
    if suffix in {".txt", ".md"}:
        return load_text(path)
    return []


def iter_supported_files(directory: Path) -> Iterable[Path]:
    for path in directory.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path
